"""마크다운 → pptx / docx / html 렌더링.

pandoc 이 있으면 우선 사용, 없으면 형식별 순수 파이썬 백엔드로 폴백한다.
pptx/docx/html 백엔드는 선택 의존성(`pip install "gigachanie[docs]"`).
"""

from __future__ import annotations

import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path

_SUPPORTED = {".pptx", ".docx", ".html", ".htm", ".md", ".pdf"}


class RenderError(RuntimeError):
    pass


# ------------------------------------------------------------ 마크다운 파싱


@dataclass
class Block:
    kind: str  # h1 | h2 | h3 | bullet | para | code
    text: str
    level: int = 0


@dataclass
class Doc:
    title: str = ""
    blocks: list[Block] = field(default_factory=list)


def parse_markdown(md: str) -> Doc:
    doc = Doc()
    in_code = False
    code_buf: list[str] = []
    for raw in md.splitlines():
        line = raw.rstrip()
        if line.startswith("```"):
            if in_code:
                doc.blocks.append(Block("code", "\n".join(code_buf)))
                code_buf = []
            in_code = not in_code
            continue
        if in_code:
            code_buf.append(raw)
            continue
        s = line.strip()
        if not s:
            continue
        if s.startswith("# "):
            t = s[2:].strip()
            if not doc.title:
                doc.title = t
            doc.blocks.append(Block("h1", t))
        elif s.startswith("## "):
            doc.blocks.append(Block("h2", s[3:].strip()))
        elif s.startswith("### "):
            doc.blocks.append(Block("h3", s[4:].strip()))
        elif s.startswith(("- ", "* ", "+ ")):
            indent = (len(line) - len(line.lstrip())) // 2
            doc.blocks.append(Block("bullet", s[2:].strip(), level=indent))
        else:
            doc.blocks.append(Block("para", s))
    if in_code and code_buf:
        doc.blocks.append(Block("code", "\n".join(code_buf)))
    return doc


# ------------------------------------------------------------ 백엔드


def _pandoc(src: Path, out: Path) -> None:
    command = ["pandoc", str(src), "-o", str(out)]
    if out.suffix.lower() in (".html", ".htm"):
        command.append("--standalone")
    proc = subprocess.run(
        command,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        timeout=60,
        check=False,
    )
    if proc.returncode != 0:
        raise RenderError(f"pandoc 실패: {proc.stderr.strip()[:300]}")


def _to_pptx(doc: Doc, out: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ModuleNotFoundError as exc:
        raise RenderError(
            "pptx 렌더에는 python-pptx 가 필요합니다: pip install \"gigachanie[docs]\""
        ) from exc

    prs = Presentation()
    blank, title_only = prs.slide_layouts[6], prs.slide_layouts[5]

    # 제목 슬라이드
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = doc.title or "Untitled"

    cur = None
    body_lines: list[tuple[int, str]] = []

    def flush() -> None:
        nonlocal cur, body_lines
        if cur is None:
            return
        slide = prs.slides.add_slide(title_only if body_lines else blank)
        if slide.shapes.title is not None:
            slide.shapes.title.text = cur
        if body_lines:
            box = slide.shapes.add_textbox(Pt(50), Pt(110), Pt(620), Pt(380))
            tf = box.text_frame
            tf.word_wrap = True
            for i, (lvl, txt) in enumerate(body_lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = ("• " if lvl >= 0 else "") + txt
                p.level = min(max(lvl, 0), 4)
                p.font.size = Pt(18)
        cur, body_lines = None, []

    for b in doc.blocks:
        if b.kind == "h1":
            continue
        if b.kind in ("h2", "h3"):
            flush()
            cur = b.text
        elif b.kind == "bullet":
            body_lines.append((b.level, b.text))
        elif b.kind == "para" or b.kind == "code":
            body_lines.append((-1, b.text))
    flush()
    prs.save(str(out))


def _to_docx(doc: Doc, out: Path) -> None:
    try:
        import docx
    except ModuleNotFoundError as exc:
        raise RenderError(
            "docx 렌더에는 python-docx 가 필요합니다: pip install \"gigachanie[docs]\""
        ) from exc

    d = docx.Document()
    for b in doc.blocks:
        if b.kind == "h1":
            d.add_heading(b.text, level=0)
        elif b.kind == "h2":
            d.add_heading(b.text, level=1)
        elif b.kind == "h3":
            d.add_heading(b.text, level=2)
        elif b.kind == "bullet":
            d.add_paragraph(b.text, style="List Bullet")
        elif b.kind == "code":
            p = d.add_paragraph()
            run = p.add_run(b.text)
            run.font.name = "Consolas"
        else:
            d.add_paragraph(b.text)
    d.save(str(out))


def _to_html(md: str, out: Path) -> None:
    try:
        import markdown as md_lib

        body = md_lib.markdown(md, extensions=["fenced_code", "tables", "toc"])
    except ModuleNotFoundError:
        body = "<pre>" + md.replace("&", "&amp;").replace("<", "&lt;") + "</pre>"
    html = (
        "<!doctype html><meta charset='utf-8'>"
        "<style>body{max-width:44rem;margin:2rem auto;padding:0 1rem;"
        "font:16px/1.7 system-ui,sans-serif}pre{background:#f4f4f4;padding:1rem;"
        "overflow:auto;border-radius:6px}code{font-family:ui-monospace,monospace}"
        "h1,h2,h3{line-height:1.25}</style>\n" + body
    )
    out.write_text(html, encoding="utf-8")


# ------------------------------------------------------------ 공개 API


def render(md_text: str, out_path: Path, *, prefer_pandoc: bool = True) -> Path:
    ext = out_path.suffix.lower()
    if ext not in _SUPPORTED:
        raise RenderError(
            f"지원하지 않는 형식: {ext} (지원: {', '.join(sorted(_SUPPORTED))})"
        )
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if ext == ".md":
        out_path.write_text(md_text, encoding="utf-8")
        return out_path

    if prefer_pandoc and shutil.which("pandoc"):
        src = out_path.with_suffix(".src.md")
        src.write_text(md_text, encoding="utf-8")
        try:
            _pandoc(src, out_path)
            return out_path
        except RenderError:
            pass
        finally:
            src.unlink(missing_ok=True)

    if ext == ".pptx":
        _to_pptx(parse_markdown(md_text), out_path)
    elif ext == ".docx":
        _to_docx(parse_markdown(md_text), out_path)
    elif ext in (".html", ".htm"):
        _to_html(md_text, out_path)
    elif ext == ".pdf":
        raise RenderError("PDF 는 pandoc(+LaTeX) 이 있어야 합니다. HTML 로 만든 뒤 인쇄하세요.")
    return out_path
